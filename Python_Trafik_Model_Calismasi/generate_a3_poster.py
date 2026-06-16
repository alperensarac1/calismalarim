from pathlib import Path
import shutil
import subprocess
import sys

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# ============================================================
# 1) PROJE AYARLARI
# ============================================================

BASE_DIR = Path(__file__).resolve().parent

DATA_PATH = BASE_DIR / "data" / "processed" / "istanbul_ilce_gunluk_trafik_model_ready.csv"
MODEL_RESULTS_PATH = BASE_DIR / "outputs" / "model_results.csv"
FEATURE_IMPORTANCE_PATH = BASE_DIR / "outputs" / "feature_importance.csv"
PREDICTIONS_PATH = BASE_DIR / "outputs" / "predictions.csv"

TEMPLATE_PATH = BASE_DIR / "NKU_BMU_Poster_Sablonu_A3_Dikey.pptx"

OUTPUT_DIR = BASE_DIR / "outputs"
OUTPUT_DIR.mkdir(exist_ok=True)

POSTER_PPTX_PATH = OUTPUT_DIR / "trafik_tahmin_a3_poster.pptx"
POSTER_PDF_PATH = OUTPUT_DIR / "trafik_tahmin_a3_poster.pdf"

CHART_DIR = OUTPUT_DIR / "poster_charts"
CHART_DIR.mkdir(exist_ok=True)


# ============================================================
# 2) POSTER İÇERİĞİ - BURAYI KENDİNE GÖRE DÜZENLEYEBİLİRSİN
# ============================================================

POSTER_TITLE = "İlçe Bazlı Trafik Tahmin ve Yol Çalışması Karar Destek Sistemi"

AUTHOR_LINE = "Alperen SARAÇ"
AFFILIATION_LINE = "Tekirdağ Namık Kemal Üniversitesi - Çorlu Mühendislik Fakültesi - Bilgisayar Mühendisliği Bölümü"
ADVISOR_LINE = "Danışman Öğretim Üyesi: ................................"

CONTACT_LINE = "📧 e-posta@nku.edu.tr  |  📍 Çorlu Mühendislik Fakültesi, Tekirdağ"

GITHUB_OR_PROJECT_URL = "https://github.com/keksteam"


# ============================================================
# 3) RENKLER
# ============================================================

COLOR_NAVY = RGBColor(20, 31, 56)
COLOR_BLUE = RGBColor(37, 99, 235)
COLOR_LIGHT_BLUE = RGBColor(239, 246, 255)
COLOR_GREEN = RGBColor(22, 163, 74)
COLOR_LIGHT_GREEN = RGBColor(240, 253, 244)
COLOR_RED = RGBColor(220, 38, 38)
COLOR_LIGHT_RED = RGBColor(254, 242, 242)
COLOR_YELLOW = RGBColor(245, 158, 11)
COLOR_LIGHT_YELLOW = RGBColor(255, 251, 235)
COLOR_GRAY = RGBColor(75, 85, 99)
COLOR_LIGHT_GRAY = RGBColor(243, 244, 246)
COLOR_BORDER = RGBColor(209, 213, 219)
COLOR_WHITE = RGBColor(255, 255, 255)


# ============================================================
# 4) VERİ OKUMA
# ============================================================

def safe_read_csv(path: Path):
    if path.exists():
        return pd.read_csv(path)

    print(f"Uyarı: Dosya bulunamadı: {path}")
    return None


def load_project_data():
    model_df = safe_read_csv(DATA_PATH)
    model_results_df = safe_read_csv(MODEL_RESULTS_PATH)
    feature_importance_df = safe_read_csv(FEATURE_IMPORTANCE_PATH)
    predictions_df = safe_read_csv(PREDICTIONS_PATH)

    if model_df is not None and "date" in model_df.columns:
        model_df["date"] = pd.to_datetime(model_df["date"], errors="coerce")

    if predictions_df is not None and "date" in predictions_df.columns:
        predictions_df["date"] = pd.to_datetime(predictions_df["date"], errors="coerce")

    return model_df, model_results_df, feature_importance_df, predictions_df


# ============================================================
# 5) GRAFİK ÜRETİMİ
# ============================================================

def save_model_comparison_chart(model_results_df):
    output_path = CHART_DIR / "model_comparison_mae.png"

    if model_results_df is None or model_results_df.empty:
        return None

    if "model" not in model_results_df.columns or "mae" not in model_results_df.columns:
        return None

    df = model_results_df.sort_values("mae").copy()

    plt.figure(figsize=(7.5, 4.2))
    plt.bar(df["model"], df["mae"])
    plt.title("Model Karşılaştırması - MAE")
    plt.xlabel("Model")
    plt.ylabel("MAE")
    plt.xticks(rotation=20, ha="right")
    plt.tight_layout()
    plt.savefig(output_path, dpi=180)
    plt.close()

    return output_path


def save_feature_importance_chart(feature_importance_df):
    output_path = CHART_DIR / "feature_importance_top10.png"

    if feature_importance_df is None or feature_importance_df.empty:
        return None

    if "feature" not in feature_importance_df.columns or "importance" not in feature_importance_df.columns:
        return None

    df = feature_importance_df.sort_values("importance", ascending=False).head(10)

    plt.figure(figsize=(7.5, 4.5))
    plt.barh(df["feature"][::-1], df["importance"][::-1])
    plt.title("En Etkili 10 Özellik")
    plt.xlabel("Önem Skoru")
    plt.tight_layout()
    plt.savefig(output_path, dpi=180)
    plt.close()

    return output_path


def save_district_average_chart(model_df):
    output_path = CHART_DIR / "district_average_top10.png"

    if model_df is None or model_df.empty:
        return None

    target_col = detect_target_column(model_df)

    if target_col is None or "district" not in model_df.columns:
        return None

    df = (
        model_df.groupby("district")[target_col]
        .mean()
        .sort_values(ascending=False)
        .head(10)
        .reset_index()
    )

    plt.figure(figsize=(7.5, 4.5))
    plt.barh(df["district"][::-1], df[target_col][::-1])
    plt.title("Ortalama Trafik Skoru En Yüksek 10 İlçe")
    plt.xlabel("Ortalama Trafik Skoru")
    plt.tight_layout()
    plt.savefig(output_path, dpi=180)
    plt.close()

    return output_path


def save_actual_predicted_chart(predictions_df):
    output_path = CHART_DIR / "actual_vs_predicted_sample.png"

    if predictions_df is None or predictions_df.empty:
        return None

    required_cols = {"actual", "predicted"}

    if not required_cols.issubset(set(predictions_df.columns)):
        return None

    sample_df = predictions_df.head(120).copy()

    plt.figure(figsize=(7.5, 4.2))
    plt.plot(sample_df["actual"].values, label="Gerçek")
    plt.plot(sample_df["predicted"].values, label="Tahmin")
    plt.title("Gerçek vs Tahmin - Test Örneği")
    plt.xlabel("Test örneği")
    plt.ylabel("Trafik Skoru")
    plt.legend()
    plt.tight_layout()
    plt.savefig(output_path, dpi=180)
    plt.close()

    return output_path


# ============================================================
# 6) ÖZET BİLGİLERİ HAZIRLAMA
# ============================================================

def detect_target_column(df):
    candidates = [
        "avg_congestion_score",
        "avg_traffic_density",
        "traffic_score",
        "target"
    ]

    for col in candidates:
        if col in df.columns:
            return col

    return None


def get_dataset_summary(model_df):
    if model_df is None or model_df.empty:
        return {
            "row_count": "-",
            "col_count": "-",
            "district_count": "-",
            "date_range": "-",
            "target_col": "-"
        }

    target_col = detect_target_column(model_df)

    if "date" in model_df.columns:
        date_min = model_df["date"].min()
        date_max = model_df["date"].max()

        if pd.notna(date_min) and pd.notna(date_max):
            date_range = f"{date_min.date()} - {date_max.date()}"
        else:
            date_range = "-"
    else:
        date_range = "-"

    district_count = model_df["district"].nunique() if "district" in model_df.columns else "-"

    return {
        "row_count": f"{len(model_df):,}",
        "col_count": str(model_df.shape[1]),
        "district_count": str(district_count),
        "date_range": date_range,
        "target_col": target_col or "-"
    }


def get_best_model_summary(model_results_df):
    if model_results_df is None or model_results_df.empty:
        return {
            "best_model": "-",
            "mae": "-",
            "rmse": "-",
            "r2": "-"
        }

    required_cols = {"model", "mae", "rmse", "r2"}

    if not required_cols.issubset(set(model_results_df.columns)):
        return {
            "best_model": "-",
            "mae": "-",
            "rmse": "-",
            "r2": "-"
        }

    best_row = model_results_df.sort_values("mae").iloc[0]

    return {
        "best_model": str(best_row["model"]),
        "mae": f"{float(best_row['mae']):.4f}",
        "rmse": f"{float(best_row['rmse']):.4f}",
        "r2": f"{float(best_row['r2']):.4f}"
    }


def get_top_features(feature_importance_df, n=5):
    if feature_importance_df is None or feature_importance_df.empty:
        return []

    if "feature" not in feature_importance_df.columns or "importance" not in feature_importance_df.columns:
        return []

    df = feature_importance_df.sort_values("importance", ascending=False).head(n)

    result = []

    for _, row in df.iterrows():
        result.append((str(row["feature"]), float(row["importance"])))

    return result


def get_model_results_table(model_results_df):
    if model_results_df is None or model_results_df.empty:
        return [["Model", "MAE", "RMSE", "R2"], ["-", "-", "-", "-"]]

    required_cols = {"model", "mae", "rmse", "r2"}

    if not required_cols.issubset(set(model_results_df.columns)):
        return [["Model", "MAE", "RMSE", "R2"], ["-", "-", "-", "-"]]

    table = [["Model", "MAE", "RMSE", "R2"]]

    df = model_results_df.sort_values("mae")

    for _, row in df.iterrows():
        table.append([
            str(row["model"]),
            f"{float(row['mae']):.4f}",
            f"{float(row['rmse']):.4f}",
            f"{float(row['r2']):.4f}"
        ])

    return table


# ============================================================
# 7) PPTX YARDIMCI FONKSİYONLARI
# ============================================================

def set_slide_size(prs):
    """
    Şablon varsa onun sayfa boyutunu kullanır.
    Şablon yoksa A3 dikey ölçü verir.
    """

    if TEMPLATE_PATH.exists():
        template = Presentation(str(TEMPLATE_PATH))
        prs.slide_width = template.slide_width
        prs.slide_height = template.slide_height
    else:
        prs.slide_width = Inches(11.69)
        prs.slide_height = Inches(16.54)


def clear_slide(slide):
    """
    Boş slide oluşturulduğu için genelde gerekmez.
    Ama şablondan gelirse şekilleri temizlemek için kullanılabilir.
    """

    shapes = list(slide.shapes)

    for shape in shapes:
        element = shape._element
        element.getparent().remove(element)


def add_textbox(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size=12,
    bold=False,
    color=COLOR_NAVY,
    align=PP_ALIGN.LEFT,
    fill_color=None,
    border_color=None,
    radius=False
):
    if radius:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )

    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color

    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = Inches(0.08)
    text_frame.margin_right = Inches(0.08)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align

    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

    return shape


def add_header(slide):
    add_textbox(
        slide,
        "T.C. TEKİRDAĞ NAMIK KEMAL ÜNİVERSİTESİ",
        x=1.65,
        y=0.25,
        w=8.4,
        h=0.28,
        font_size=11,
        bold=True,
        color=COLOR_NAVY,
        align=PP_ALIGN.CENTER
    )

    add_textbox(
        slide,
        "Çorlu Mühendislik Fakültesi - Bilgisayar Mühendisliği Bölümü",
        x=1.65,
        y=0.55,
        w=8.4,
        h=0.24,
        font_size=9,
        bold=False,
        color=COLOR_GRAY,
        align=PP_ALIGN.CENTER
    )

    add_textbox(
        slide,
        POSTER_TITLE,
        x=1.1,
        y=0.95,
        w=9.5,
        h=0.72,
        font_size=20,
        bold=True,
        color=COLOR_NAVY,
        align=PP_ALIGN.CENTER
    )

    add_textbox(
        slide,
        AUTHOR_LINE,
        x=1.1,
        y=1.75,
        w=9.5,
        h=0.25,
        font_size=10,
        bold=True,
        color=COLOR_BLUE,
        align=PP_ALIGN.CENTER
    )

    add_textbox(
        slide,
        AFFILIATION_LINE,
        x=1.1,
        y=2.05,
        w=9.5,
        h=0.24,
        font_size=8.5,
        color=COLOR_GRAY,
        align=PP_ALIGN.CENTER
    )

    add_textbox(
        slide,
        ADVISOR_LINE,
        x=1.1,
        y=2.32,
        w=9.5,
        h=0.22,
        font_size=8.5,
        color=COLOR_GRAY,
        align=PP_ALIGN.CENTER
    )

    # Logo alanı
    add_textbox(
        slide,
        "ÜNİVERSİTE\nLOGOSU",
        x=0.25,
        y=0.25,
        w=1.15,
        h=1.0,
        font_size=8,
        bold=True,
        color=COLOR_GRAY,
        align=PP_ALIGN.CENTER,
        fill_color=COLOR_LIGHT_GRAY,
        border_color=COLOR_BORDER,
        radius=True
    )

    # QR alanı
    add_textbox(
        slide,
        "QR KOD\nGitHub / Web",
        x=10.25,
        y=0.25,
        w=1.15,
        h=1.0,
        font_size=7.5,
        bold=True,
        color=COLOR_GRAY,
        align=PP_ALIGN.CENTER,
        fill_color=COLOR_LIGHT_GRAY,
        border_color=COLOR_BORDER,
        radius=True
    )


def add_section_box(slide, title, body, x, y, w, h, fill_color=COLOR_WHITE):
    slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )

    box = slide.shapes[-1]
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = COLOR_BORDER
    box.line.width = Pt(0.75)

    add_textbox(
        slide,
        title,
        x=x + 0.12,
        y=y + 0.08,
        w=w - 0.24,
        h=0.28,
        font_size=11,
        bold=True,
        color=COLOR_BLUE
    )

    add_textbox(
        slide,
        body,
        x=x + 0.12,
        y=y + 0.42,
        w=w - 0.24,
        h=h - 0.52,
        font_size=7.7,
        color=COLOR_NAVY
    )


def add_metric_card(slide, title, value, x, y, w, h, color):
    add_textbox(
        slide,
        title,
        x=x,
        y=y,
        w=w,
        h=0.25,
        font_size=7.5,
        bold=True,
        color=COLOR_GRAY,
        align=PP_ALIGN.CENTER,
        fill_color=COLOR_LIGHT_GRAY,
        border_color=COLOR_BORDER,
        radius=True
    )

    add_textbox(
        slide,
        value,
        x=x,
        y=y + 0.28,
        w=w,
        h=h - 0.28,
        font_size=13,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        fill_color=COLOR_WHITE,
        border_color=COLOR_BORDER,
        radius=True
    )


def add_table(slide, table_data, x, y, w, h, font_size=6.5):
    rows = len(table_data)
    cols = len(table_data[0])

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )

    table = table_shape.table

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(table_data[r][c])

            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY if r == 0 else COLOR_WHITE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER

                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = True if r == 0 else False
                    run.font.color.rgb = COLOR_NAVY

    return table_shape


def add_picture_if_exists(slide, image_path, x, y, w, h):
    if image_path is None:
        return

    image_path = Path(image_path)

    if not image_path.exists():
        return

    slide.shapes.add_picture(
        str(image_path),
        Inches(x),
        Inches(y),
        width=Inches(w),
        height=Inches(h)
    )


def add_footer(slide):
    add_textbox(
        slide,
        CONTACT_LINE,
        x=0.35,
        y=16.05,
        w=10.95,
        h=0.26,
        font_size=7.5,
        color=COLOR_GRAY,
        align=PP_ALIGN.CENTER
    )


# ============================================================
# 8) POSTER OLUŞTURMA
# ============================================================

def build_poster():
    model_df, model_results_df, feature_importance_df, predictions_df = load_project_data()

    summary = get_dataset_summary(model_df)
    best_model = get_best_model_summary(model_results_df)
    top_features = get_top_features(feature_importance_df, n=5)

    model_chart = save_model_comparison_chart(model_results_df)
    feature_chart = save_feature_importance_chart(feature_importance_df)
    district_chart = save_district_average_chart(model_df)
    actual_pred_chart = save_actual_predicted_chart(predictions_df)

    prs = Presentation()
    set_slide_size(prs)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Arka plan
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    add_header(slide)

    # ------------------------------------------------------------
    # Metrik kartları
    # ------------------------------------------------------------

    add_metric_card(
        slide,
        "Veri Satırı",
        summary["row_count"],
        x=0.35,
        y=2.75,
        w=2.05,
        h=0.7,
        color=COLOR_BLUE
    )

    add_metric_card(
        slide,
        "İlçe Sayısı",
        summary["district_count"],
        x=2.55,
        y=2.75,
        w=2.05,
        h=0.7,
        color=COLOR_GREEN
    )

    add_metric_card(
        slide,
        "En İyi Model",
        best_model["best_model"],
        x=4.75,
        y=2.75,
        w=2.05,
        h=0.7,
        color=COLOR_BLUE
    )

    add_metric_card(
        slide,
        "R2",
        best_model["r2"],
        x=6.95,
        y=2.75,
        w=2.05,
        h=0.7,
        color=COLOR_GREEN
    )

    add_metric_card(
        slide,
        "MAE",
        best_model["mae"],
        x=9.15,
        y=2.75,
        w=2.05,
        h=0.7,
        color=COLOR_RED
    )

    # ------------------------------------------------------------
    # Sol kolon
    # ------------------------------------------------------------

    abstract_text = (
        "Bu çalışma, İstanbul ilçeleri için günlük trafik sıkışıklığı tahmini yapan "
        "makine öğrenmesi tabanlı bir karar destek sistemi geliştirmektedir. "
        "Geçmiş trafik verilerinden tarihsel, gecikmeli ve hareketli ortalama "
        "özellikleri üretilmiş; farklı regresyon modelleri karşılaştırılmıştır. "
        "En iyi model, yol çalışması planlamasında uygun gün seçimi için "
        "senaryo tabanlı arayüzde kullanılmıştır."
    )

    add_section_box(
        slide,
        "📋 ÖZET",
        abstract_text,
        x=0.35,
        y=3.7,
        w=3.45,
        h=1.55,
        fill_color=COLOR_WHITE
    )

    motivation_text = (
        "• Şehir içi trafik yoğunluğu yol çalışması planlamasını doğrudan etkiler.\n"
        "• Yanlış zamanda yapılan çalışma trafik sıkışıklığını artırabilir.\n"
        "• Amaç, ilçe ve gün bazlı trafik tahmini üretmek ve çalışma için uygun günü belirlemektir.\n"
        "• Sistem, tahmin sonucunu Çok uygun / Uygun / Uygun değil olarak yorumlar."
    )

    add_section_box(
        slide,
        "🎯 GİRİŞ VE MOTİVASYON",
        motivation_text,
        x=0.35,
        y=5.45,
        w=3.45,
        h=1.65,
        fill_color=COLOR_LIGHT_BLUE
    )

    dataset_text = (
        f"Veri kümesi: İstanbul trafik verileri\n"
        f"Satır sayısı: {summary['row_count']}\n"
        f"Sütun sayısı: {summary['col_count']}\n"
        f"İlçe sayısı: {summary['district_count']}\n"
        f"Tarih aralığı: {summary['date_range']}\n"
        f"Hedef değişken: {summary['target_col']}\n\n"
        "Ön işlem:\n"
        "• İlçe-gün seviyesinde gruplama\n"
        "• Tarihsel özellik üretimi\n"
        "• Lag ve rolling feature üretimi\n"
        "• Veri sızıntısı oluşturabilecek sütunların çıkarılması"
    )

    add_section_box(
        slide,
        "🗄️ VERİ KÜMESİ",
        dataset_text,
        x=0.35,
        y=7.3,
        w=3.45,
        h=2.45,
        fill_color=COLOR_WHITE
    )

    features_text = "En etkili özellikler:\n"

    if top_features:
        for feature, importance in top_features:
            features_text += f"• {feature}: {importance:.4f}\n"
    else:
        features_text += "• Feature importance dosyası bulunamadı.\n"

    features_text += (
        "\nYorum:\n"
        "Modelin en fazla geçmiş günlere ait trafik davranışlarından etkilendiği görülmektedir. "
        "Bu durum trafik yoğunluğunun zamansal süreklilik taşıdığını göstermektedir."
    )

    add_section_box(
        slide,
        "⚙️ ÖZELLİK MÜHENDİSLİĞİ",
        features_text,
        x=0.35,
        y=9.95,
        w=3.45,
        h=2.2,
        fill_color=COLOR_LIGHT_GREEN
    )

    # ------------------------------------------------------------
    # Orta kolon
    # ------------------------------------------------------------

    method_text = (
        "1. Ham trafik kayıtları okundu.\n"
        "2. Tarih ve ilçe bazında günlük trafik skoru üretildi.\n"
        "3. Ay, gün, hafta sonu, resmi tatil ve döngüsel tarih özellikleri eklendi.\n"
        "4. lag_1, lag_2, lag_3, lag_7 ve rolling mean/std değişkenleri üretildi.\n"
        "5. Ridge, Random Forest, Extra Trees ve Gradient Boosting modelleri test edildi.\n"
        "6. En düşük MAE değerine sahip model seçildi.\n"
        "7. Kullanıcı senaryoları ile tahminler karar destek çıktısına dönüştürüldü."
    )

    add_section_box(
        slide,
        "🔬 YÖNTEM VE DENEYSEL KURULUM",
        method_text,
        x=4.05,
        y=3.7,
        w=3.55,
        h=2.5,
        fill_color=COLOR_WHITE
    )

    add_section_box(
        slide,
        "📈 MODEL SONUÇLARI",
        (
            f"En iyi model: {best_model['best_model']}\n"
            f"MAE: {best_model['mae']}\n"
            f"RMSE: {best_model['rmse']}\n"
            f"R2: {best_model['r2']}\n\n"
            "MAE düşük olduğunda modelin ortalama hata miktarı azalır. "
            "R2 değerinin yüksek olması, modelin hedef değişkendeki varyansı güçlü şekilde açıkladığını gösterir."
        ),
        x=4.05,
        y=6.4,
        w=3.55,
        h=1.75,
        fill_color=COLOR_LIGHT_BLUE
    )

    model_table = get_model_results_table(model_results_df)

    add_table(
        slide,
        model_table,
        x=4.15,
        y=8.35,
        w=3.35,
        h=1.25,
        font_size=5.6
    )

    if model_chart:
        add_picture_if_exists(
            slide,
            model_chart,
            x=4.05,
            y=9.85,
            w=3.55,
            h=2.05
        )

    if actual_pred_chart:
        add_picture_if_exists(
            slide,
            actual_pred_chart,
            x=4.05,
            y=12.05,
            w=3.55,
            h=2.05
        )

    # ------------------------------------------------------------
    # Sağ kolon
    # ------------------------------------------------------------

    decision_text = (
        "Senaryo parametreleri:\n"
        "• Yol çalışması var/yok\n"
        "• Kapalı şerit sayısı\n"
        "• Çalışma zamanı: Gündüz / Gece\n"
        "• Yağmur seviyesi\n"
        "• Etkinlik yoğunluğu\n"
        "• Okul günü bilgisi\n"
        "• Kaza riski\n"
        "• Toplu taşıma aksaması\n\n"
        "Çıktı:\n"
        "• Trafik risk seviyesi\n"
        "• Uygunluk: Çok uygun / Uygun / Uygun değil\n"
        "• Gerekçe ve öneri metni\n"
        "• Alternatif gün önerileri"
    )

    add_section_box(
        slide,
        "🧠 KARAR DESTEK SİSTEMİ",
        decision_text,
        x=7.85,
        y=3.7,
        w=3.45,
        h=2.6,
        fill_color=COLOR_LIGHT_YELLOW
    )

    if feature_chart:
        add_picture_if_exists(
            slide,
            feature_chart,
            x=7.85,
            y=6.5,
            w=3.45,
            h=2.25
        )

    if district_chart:
        add_picture_if_exists(
            slide,
            district_chart,
            x=7.85,
            y=8.95,
            w=3.45,
            h=2.25
        )

    discussion_text = (
        "• Geçmiş trafik değerleri, model performansında en belirleyici değişkenlerdir.\n"
        "• Gece çalışma senaryosu trafik etkisini azaltmaktadır.\n"
        "• Yüksek riskli günlerde uygunluk en fazla 'Uygun' veya 'Uygun değil' olarak değerlendirilmelidir.\n"
        "• Sistem, yalnızca tahmin üretmek yerine karar gerekçesi ve öneri de sunmaktadır."
    )

    add_section_box(
        slide,
        "💬 TARTIŞMA",
        discussion_text,
        x=7.85,
        y=11.4,
        w=3.45,
        h=1.6,
        fill_color=COLOR_WHITE
    )

    conclusion_text = (
        "Bu çalışmada İstanbul ilçeleri için günlük trafik tahmini yapan ve yol çalışması planlamasına "
        "destek veren bir sistem geliştirilmiştir. Model sonuçları, geçmiş trafik davranışlarının güçlü "
        "bir tahmin değişkeni olduğunu göstermektedir.\n\n"
        "Gelecek çalışmalar:\n"
        "• Hava durumu verilerinin eklenmesi\n"
        "• Gerçek kaza ve etkinlik kayıtlarının modele dahil edilmesi\n"
        "• Harita tabanlı risk görselleştirmesi\n"
        "• XGBoost / LightGBM / CatBoost modelleriyle karşılaştırma"
    )

    add_section_box(
        slide,
        "✅ SONUÇ VE GELECEK ÇALIŞMA",
        conclusion_text,
        x=0.35,
        y=14.25,
        w=10.95,
        h=1.55,
        fill_color=COLOR_LIGHT_GREEN
    )

    add_footer(slide)

    prs.save(str(POSTER_PPTX_PATH))

    print("\nA3 poster başarıyla oluşturuldu:")
    print(POSTER_PPTX_PATH)

    try_convert_to_pdf(POSTER_PPTX_PATH, POSTER_PDF_PATH)


# ============================================================
# 9) PDF'E ÇEVİRME - OPSİYONEL
# ============================================================

def try_convert_to_pdf(pptx_path, pdf_path):
    """
    Bilgisayarda LibreOffice varsa otomatik PDF üretir.
    Yoksa sadece PPTX oluşturur.
    """

    libreoffice_commands = [
        "soffice",
        "libreoffice"
    ]

    selected_command = None

    for command in libreoffice_commands:
        if shutil.which(command):
            selected_command = command
            break

    if selected_command is None:
        print("\nPDF üretimi atlandı.")
        print("PDF için LibreOffice kurulu olmalı veya PPTX'i PowerPoint üzerinden PDF olarak kaydedebilirsin.")
        return

    try:
        subprocess.run(
            [
                selected_command,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(pdf_path.parent),
                str(pptx_path)
            ],
            check=True
        )

        print("\nPDF de oluşturuldu:")
        print(pdf_path)

    except Exception as e:
        print("\nPDF dönüştürme sırasında hata oluştu.")
        print(e)


# ============================================================
# 10) ÇALIŞTIRMA
# ============================================================

if __name__ == "__main__":
    build_poster()